"""
Export Google Docs to Markdown, Slides to Markdown (via PPTX),
or Sheets to CSV.

Requires gcloud CLI and python-pptx (for Slides export).

python gdoc2md.py <google-doc-or-slides-or-sheets-url> [output]
"""

import json
import re
import subprocess
import sys
import time
from io import BytesIO
from pathlib import Path
from urllib.request import Request, urlopen
from urllib.error import HTTPError


# tolerates trailing segments like /edit, /view, ?usp=sharing
VALID_URL_RE = re.compile(
    r"^https://docs\.google\.com/"
    r"(?P<mode>document|presentation|spreadsheets)/d/(?P<id>[a-zA-Z0-9_-]+)"
)

MODE_MAP = {
    "document": "doc",
    "presentation": "slides",
    "spreadsheets": "sheets",
}

EXTENSIONS = {"doc": ".md", "slides": ".md", "sheets": ".csv"}


# ---------------------------------------------------------------------------
# Argument parsing & validation
# ---------------------------------------------------------------------------

def parse_and_validate_args():
    if len(sys.argv) < 2:
        print(
            f"Usage: {sys.argv[0]} "
            "<google-doc-or-slides-or-sheets-url> [output]"
        )
        sys.exit(1)

    url = sys.argv[1]
    match = VALID_URL_RE.match(url)
    if not match:
        print(
            "Error: URL must be a Google Docs, Slides, or Sheets URL "
            "(https://docs.google.com/...)",
            file=sys.stderr,
        )
        sys.exit(1)

    mode = MODE_MAP[match.group("mode")]
    file_id = match.group("id")

    explicit_output = sys.argv[2] if len(sys.argv) > 2 else None
    output = explicit_output or f"{file_id}{EXTENSIONS[mode]}"

    return file_id, output, mode


# ---------------------------------------------------------------------------
# Dependency check
# ---------------------------------------------------------------------------

def check_dependencies():
    result = subprocess.run(["gcloud", "version"], capture_output=True)
    if result.returncode != 0:
        print("Error: gcloud CLI is not installed.", file=sys.stderr)
        print(
            "  Install: "
            "https://cloud.google.com/sdk/docs/install",
            file=sys.stderr,
        )
        sys.exit(1)


# ---------------------------------------------------------------------------
# Auth — single source of truth for obtaining a token
# ---------------------------------------------------------------------------

def get_token() -> str:
    """
    Return a valid access token, prompting the user to log in if needed.
    Raises SystemExit on unrecoverable failure.
    """
    result = subprocess.run(
        ["gcloud", "auth", "print-access-token"],
        capture_output=True, text=True,
    )
    if result.returncode == 0 and result.stdout.strip():
        return result.stdout.strip()

    print("No active credentials found. Authenticating with Google...")
    login = subprocess.run(
        ["gcloud", "auth", "login", "--enable-gdrive-access"],
    )
    if login.returncode != 0:
        print(
            "Error: Authentication was cancelled or failed.",
            file=sys.stderr,
        )
        sys.exit(1)

    # Re-fetch after successful login
    result = subprocess.run(
        ["gcloud", "auth", "print-access-token"],
        capture_output=True, text=True,
    )
    if result.returncode != 0 or not result.stdout.strip():
        print(
            "Error: Could not obtain access token after login.",
            file=sys.stderr,
        )
        sys.exit(1)

    return result.stdout.strip()


# ---------------------------------------------------------------------------
# Download
# ---------------------------------------------------------------------------

def download(url: str, token: str, retries: int = 3) -> bytes:
    req = Request(url, headers={"Authorization": f"Bearer {token}"})
    for attempt in range(retries + 1):
        try:
            with urlopen(req) as resp:
                return resp.read()
        except HTTPError as e:
            if e.code == 429 and attempt < retries:
                wait = 2 ** attempt
                print(
                    f"Rate limited (429), retrying in {wait}s...",
                    file=sys.stderr,
                )
                time.sleep(wait)
                continue
            messages = {
                401: (
                    "Authentication failed (401). "
                    "Try: gcloud auth login "
                    "--enable-gdrive-access"
                ),
                403: (
                    "Access denied (403). Check you have "
                    "permission to access this file."
                ),
                404: "Not found (404). Check the URL is correct.",
            }
            print(
                f"Error: {messages.get(e.code, f'HTTP {e.code}')}",
                file=sys.stderr,
            )
            sys.exit(1)


# ---------------------------------------------------------------------------
# PPTX → Markdown conversion
# ---------------------------------------------------------------------------

def pptx_to_markdown(data: bytes) -> str:
    """
    Convert PPTX bytes to structured Markdown.

    Note: images, charts, and other non-text/non-table shapes are not exported.
    Only text frames and tables are extracted.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print(
            "Error: python-pptx is required for Slides export.",
            file=sys.stderr,
        )
        print("  Install: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(BytesIO(data))
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        lines.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.replace("\x0b", "\n").strip()
                    if not text:
                        continue
                    if paragraph.level > 0:
                        indent = "  " * (paragraph.level - 1)
                        for subline in text.split("\n"):
                            lines.append(f"{indent}- {subline}")
                    else:
                        lines.append(text)
                lines.append("")

            elif shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    cells = [
                        cell.text.strip().replace("|", "\\|")
                        for cell in row.cells
                    ]
                    lines.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        sep = "| " + " | ".join(
                            ["---"] * len(cells)
                        ) + " |"
                        lines.append(sep)
                lines.append("")

        # Images, charts, and other shape types are not exported
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("> **Notes:** " + notes_text)
                lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Sheets metadata
# ---------------------------------------------------------------------------

def get_sheet_metadata(file_id: str, token: str):
    """Return list of (gid, title) for each sheet."""
    api_url = (
        "https://sheets.googleapis.com/v4/spreadsheets/"
        f"{file_id}?fields=sheets.properties"
    )
    data = download(api_url, token)
    info = json.loads(data)
    return [
        (
            s["properties"]["sheetId"],
            s["properties"]["title"],
        )
        for s in info["sheets"]
    ]


def _sanitize_filename(name: str) -> str:
    """Replace filesystem-unsafe characters with underscores."""
    return re.sub(r'[\\/*?:"<>|]', "_", name)


# ---------------------------------------------------------------------------
# Fetch & write
# ---------------------------------------------------------------------------

def fetch(file_id: str, output: str, mode: str):
    token = get_token()
    base = "https://docs.google.com"

    if mode == "sheets":
        _fetch_sheets(file_id, output, token, base)
        return

    export_urls = {
        "slides": (
            f"{base}/presentation/d/{file_id}"
            "/export?format=pptx"
        ),
        "doc": (
            f"{base}/document/d/{file_id}"
            "/export?format=md"
        ),
    }

    data = download(export_urls[mode], token)
    output_path = Path(output)

    if output_path.exists():
        print(
            f"Warning: overwriting existing file '{output}'",
            file=sys.stderr,
        )

    if mode == "slides":
        output_path.write_text(
            pptx_to_markdown(data), encoding="utf-8"
        )
    else:
        output_path.write_bytes(data)

    print(f"Saved to {output}")


def _fetch_sheets(
    file_id: str, output: str, token: str, base: str
):
    """Export every sheet in a spreadsheet as a separate CSV."""
    try:
        sheets = get_sheet_metadata(file_id, token)
    except SystemExit:
        # Sheets API not enabled — fall back to default first sheet
        print(
            "Warning: Could not fetch sheet metadata "
            "(Sheets API may not be enabled). "
            "Exporting first sheet only.",
            file=sys.stderr,
        )
        sheets = [(0, "Sheet1")]
    out_path = Path(output)
    stem = out_path.stem
    parent = out_path.parent

    if len(sheets) == 1:
        gid, title = sheets[0]
        url = (
            f"{base}/spreadsheets/d/{file_id}"
            f"/export?format=csv&gid={gid}"
        )
        data = download(url, token)
        if out_path.exists():
            print(
                "Warning: overwriting existing "
                f"file '{out_path}'",
                file=sys.stderr,
            )
        out_path.write_bytes(data)
        print(f"Saved to {out_path}")
        return

    for gid, title in sheets:
        safe_title = _sanitize_filename(title)
        csv_path = parent / f"{stem}_{safe_title}.csv"
        url = (
            f"{base}/spreadsheets/d/{file_id}"
            f"/export?format=csv&gid={gid}"
        )
        data = download(url, token)
        if csv_path.exists():
            print(
                "Warning: overwriting existing "
                f"file '{csv_path}'",
                file=sys.stderr,
            )
        csv_path.write_bytes(data)
        print(f"Saved to {csv_path}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    # Validate args first — fast failure before any subprocess calls
    file_id, output, mode = parse_and_validate_args()
    check_dependencies()
    fetch(file_id, output, mode)


if __name__ == "__main__":
    main()